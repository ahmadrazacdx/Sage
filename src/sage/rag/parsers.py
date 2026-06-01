"""
Format parsers for Sage's preprocessing pipeline.

Available Parsers:
  .pdf  : PyMuPDF (primary) | pdfplumber fallback | OCR fallback
  .pptx : python-pptx: shapes, notes, tables
  .docx : python-docx: paragraphs, tables, headings
  .md   : raw UTF-8 read
  .txt  : raw UTF-8 / chardet fallback

Usage:
    from sage.rag.parsers import dispatch_parser, ParsedDocument
    doc = dispatch_parser(meta)   # meta: DocumentMetadata
"""

from __future__ import annotations

import re
import warnings
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

import structlog

from sage.rag.corpus import DocumentMetadata

log = structlog.get_logger(__name__)
warnings.filterwarnings("ignore", category=DeprecationWarning, module="pdf2image")


_PDF_MAGIC: bytes = b"%PDF"
_OCR_MAX_ZERO_TEXT_RATIO: float = 0.30
# OCR ligature repairs
_OCR_LIGATURE_MAP: dict[str, str] = {
    "ﬁ": "fi",
    "ﬀ": "ff",
    "ﬂ": "fl",
    "ﬃ": "ffi",
    "ﬄ": "ffl",
    "ﬅ": "ft",
    "ﬆ": "st",
}

@dataclass
class PageBlock:
    """Text extracted from a single page or slide, with its 1-based number."""

    page_number: int
    text: str


@dataclass
class ParsedDocument:
    """
    Raw extraction result for one source document.

    Attributes:
        blocks:       Ordered list of page/slide text blocks.
        page_count:   Total pages or slides in the document.
        ocr_applied:  True if OCR was invoked for any page.
        ocr_engine:   "tesseract" | "easyocr" | "" when not applied.
        ocr_failed:   True when OCR was required but failed and no usable text was produced.
    """

    blocks: list[PageBlock] = field(default_factory=list)
    page_count: int = 0
    ocr_applied: bool = False
    ocr_engine: str = ""
    ocr_failed: bool = False

    def full_text(self) -> str:
        """Return all block texts joined by double-newline."""
        return "\n\n".join(b.text for b in self.blocks if b.text.strip())

def _require(pkg: str, install: str) -> None:
    """Raise ImportError with an actionable message if pkg is missing."""
    raise ImportError(
        f"Required package '{pkg}' is not installed. "
        f"Install it with:  uv add {install}"
    )


def _validate_pdf_magic(path: Path) -> bool:
    """Return True if the file starts with the PDF magic bytes '%PDF'."""
    try:
        with path.open("rb") as fh:
            return fh.read(4) == _PDF_MAGIC
    except OSError:
        return False

# PDF Parser
def _parse_pdf_pymupdf(path: Path) -> list[PageBlock]:
    """Primary PDF parser using PyMuPDF (fitz)."""
    try:
        import fitz  # type: ignore[import-untyped]  # pymupdf
    except ImportError:
        _require("pymupdf", "pymupdf")

    blocks: list[PageBlock] = []
    doc = fitz.open(str(path))
    try:
        page_count = doc.page_count
        for i in range(page_count):
            page = doc[i]
            try:
                text: str = page.get_text("text") or ""
            finally:
                page = None  # type: ignore[assignment]
            blocks.append(PageBlock(page_number=i + 1, text=text))
    finally:
        doc.close()
    return blocks

def _parse_pdf_pdfplumber(path: Path) -> list[PageBlock]:
    """PDF parser using pdfplumber with layout=True."""
    try:
        import pdfplumber  # type: ignore[import-untyped]
    except ImportError:
        _require("pdfplumber", "pdfplumber")

    blocks: list[PageBlock] = []
    with pdfplumber.open(str(path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text: str = page.extract_text(layout=True) or ""
            blocks.append(PageBlock(page_number=i, text=text))
    return blocks


def _needs_ocr(blocks: list[PageBlock], min_chars: int) -> bool:
    """Decide whether OCR is needed based on extraction quality."""
    if not blocks:
        return True
    total_chars = sum(len(b.text) for b in blocks)
    zero_pages = sum(1 for b in blocks if not b.text.strip())
    avg_chars = total_chars / len(blocks)
    zero_ratio = zero_pages / len(blocks)
    return avg_chars < min_chars or zero_ratio > _OCR_MAX_ZERO_TEXT_RATIO


def _apply_ocr_ligature_cleanup(text: str) -> str:
    """Fix common OCR ligature errors and hyphenation artifacts."""
    text = re.sub(r"-\n(\w)", r"\1", text)
    text = re.sub(r"(?<!\n)\n(?!\n)", " ", text)
    for ligature, replacement in _OCR_LIGATURE_MAP.items():
        text = text.replace(ligature, replacement)
    return text


def _parse_pdf_ocr_tesseract(
    path: Path, dpi: int, language: str
) -> tuple[list[PageBlock], str]:
    """OCR via pytesseract + pdf2image. Returns (blocks, engine_name)."""
    try:
        import pytesseract  # type: ignore[import-untyped]
        from pdf2image import convert_from_path  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            f"OCR requires pytesseract and pdf2image: uv add pytesseract pdf2image. "
            f"Also install Tesseract binary and Poppler. Error: {exc}"
        ) from exc

    images = convert_from_path(str(path), dpi=dpi)
    blocks: list[PageBlock] = []
    for i, img in enumerate(images, start=1):
        text: str = pytesseract.image_to_string(img, lang=language) or ""
        text = _apply_ocr_ligature_cleanup(text)
        blocks.append(PageBlock(page_number=i, text=text))
    return blocks, "tesseract"


def _parse_pdf_ocr_easyocr(path: Path, dpi: int) -> tuple[list[PageBlock], str]:
    """OCR fallback via easyocr. Returns (blocks, engine_name)."""
    try:
        import easyocr  # type: ignore[import-untyped]
        from pdf2image import convert_from_path  # type: ignore[import-untyped]
    except ImportError as exc:
        raise ImportError(
            f"OCR fallback requires easyocr and pdf2image: uv add easyocr pdf2image. "
            f"Error: {exc}"
        ) from exc

    reader = easyocr.Reader(["en"], gpu=False)
    images = convert_from_path(str(path), dpi=dpi)
    blocks: list[PageBlock] = []
    for i, img in enumerate(images, start=1):
        import numpy as np  # type: ignore[import-untyped]

        results = reader.readtext(np.array(img), detail=0)
        text = _apply_ocr_ligature_cleanup(" ".join(results))
        blocks.append(PageBlock(page_number=i, text=text))
    return blocks, "easyocr"


def parse_pdf(
    path: Path,
    meta: DocumentMetadata,
    min_chars_per_page: int,
    ocr_engine: str,
    ocr_dpi: int,
    ocr_language: str,
) -> ParsedDocument:
    """
    Parse a PDF file through the full extraction cascade:
    pdfplumber → PyMuPDF fallback → OCR fallback.

    Mutates *meta* to set ``ocr_applied`` and ``ocr_engine`` when OCR fires.
    """
    if not _validate_pdf_magic(path):
        log.warning(
            "pdf_invalid_magic",
            path=str(path),
            hint="File does not start with %PDF — may be corrupt or misnamed.",
        )

    # Primary extraction
    blocks: list[PageBlock] = []
    try:
        blocks = _parse_pdf_pymupdf(path)
        log.debug("pdf_pymupdf_ok", path=str(path), pages=len(blocks))
    except Exception as exc:
        log.warning(
            "pdf_pymupdf_failed",
            path=str(path),
            exc_type=type(exc).__name__,
            error=str(exc)[:200],
            hint="Falling back to PyMuPDF.",
        )

    # PyMuPDF fallback when pdfplumber yields < 50 chars/page on average
    if not blocks or _needs_ocr(blocks, min_chars_per_page):
        try:
            fallback_blocks = _parse_pdf_pdfplumber(path)
            if fallback_blocks and not _needs_ocr(fallback_blocks, min_chars_per_page):
                blocks = fallback_blocks
                log.debug("pdf_pdfplumber_fallback_used", path=str(path))
        except Exception as exc:
            log.warning(
                "pdf_pdfplumber_failed",
                path=str(path),
                exc_type=type(exc).__name__,
                error=str(exc)[:200],
            )

    # OCR fallback for scanned PDFs
    ocr_applied = False
    ocr_failed = False
    actual_ocr_engine = ""
    if not blocks or _needs_ocr(blocks, min_chars_per_page):
        log.info(
            "pdf_ocr_triggered",
            path=str(path),
            engine=ocr_engine,
            dpi=ocr_dpi,
        )
        try:
            if ocr_engine == "tesseract":
                blocks, actual_ocr_engine = _parse_pdf_ocr_tesseract(
                    path, ocr_dpi, ocr_language
                )
            else:
                blocks, actual_ocr_engine = _parse_pdf_ocr_easyocr(path, ocr_dpi)
            ocr_applied = True
        except Exception as exc:
            ocr_failed = True
            log.warning(
                "pdf_ocr_failed",
                path=str(path),
                engine=ocr_engine,
                exc_type=type(exc).__name__,
                error=str(exc)[:300],
            )

    if ocr_failed and any(block.text.strip() for block in blocks):
        ocr_failed = False

    meta.ocr_applied = ocr_applied
    meta.ocr_engine = actual_ocr_engine

    return ParsedDocument(
        blocks=blocks,
        page_count=len(blocks),
        ocr_applied=ocr_applied,
        ocr_engine=actual_ocr_engine,
        ocr_failed=ocr_failed,
    )


# ---------------------------------------------------------------------------
# PPTX parser
# ---------------------------------------------------------------------------


def parse_pptx(path: Path) -> ParsedDocument:
    """
    Extract text from all slides: title, body shapes, notes, table cells.

    Text is tagged with 1-based slide_number (stored as page_number).
    """
    try:
        from pptx import Presentation  # type: ignore[import-untyped]
        from pptx.util import Pt  # noqa: F401 — ensures pptx is importable
    except ImportError:
        _require("python-pptx", "python-pptx")

    prs = Presentation(str(path))
    blocks: list[PageBlock] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []

        for shape in slide.shapes:
            # Text frames (titles, bodies, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs if run.text)
                    if line.strip():
                        parts.append(line)

            # Tables — row-major order
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = "\t".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)

        # Slide notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes_text = notes_tf.text.strip() if notes_tf else ""
            if notes_text:
                parts.append(f"[Notes] {notes_text}")

        blocks.append(PageBlock(page_number=slide_num, text="\n".join(parts)))

    return ParsedDocument(blocks=blocks, page_count=len(blocks))


# ---------------------------------------------------------------------------
# DOCX parser
# ---------------------------------------------------------------------------


def parse_docx(path: Path) -> ParsedDocument:
    """
    Extract paragraphs, tables, headers, and footers from a DOCX file.

    Heading styles are converted to markdown # hierarchy.  All content
    is returned as a single page block (DOCX has no inherent page breaks
    in the python-docx model).
    """
    try:
        import docx as python_docx  # type: ignore[import-untyped]
    except ImportError:
        _require("python-docx", "python-docx")

    document = python_docx.Document(str(path))
    parts: list[str] = []

    _HEADING_PREFIX: dict[str, str] = {
        "Heading 1": "# ",
        "Heading 2": "## ",
        "Heading 3": "### ",
        "Heading 4": "#### ",
    }

    for para in document.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        prefix = _HEADING_PREFIX.get(para.style.name, "")
        parts.append(f"{prefix}{text}")

    for table in document.tables:
        for row in table.rows:
            row_text = "\t".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)

    full_text = "\n".join(parts)
    return ParsedDocument(
        blocks=[PageBlock(page_number=1, text=full_text)],
        page_count=1,
    )


# ---------------------------------------------------------------------------
# Plain text / Markdown parser
# ---------------------------------------------------------------------------


def parse_text(path: Path) -> ParsedDocument:
    """
    Read .md and .txt files as UTF-8, falling back to chardet on decode error.

    Content is returned as a single page block.
    """
    text: Optional[str] = None

    # Primary: UTF-8
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        log.debug("text_utf8_failed", path=str(path), hint="Trying chardet detection.")

    # Fallback: chardet encoding detection
    if text is None:
        try:
            import chardet  # type: ignore[import-untyped]
        except ImportError:
            _require("chardet", "chardet")

        raw_bytes = path.read_bytes()
        detected = chardet.detect(raw_bytes)
        encoding: str = detected.get("encoding") or "utf-8"
        log.debug(
            "text_chardet_detected",
            path=str(path),
            encoding=encoding,
            confidence=detected.get("confidence"),
        )
        text = raw_bytes.decode(encoding, errors="replace")

    return ParsedDocument(
        blocks=[PageBlock(page_number=1, text=text or "")],
        page_count=1,
    )


# ---------------------------------------------------------------------------
# Main dispatcher
# ---------------------------------------------------------------------------


def dispatch_parser(
    meta: DocumentMetadata,
    min_chars_per_page: int = 50,
    ocr_engine: str = "tesseract",
    ocr_dpi: int = 300,
    ocr_language: str = "eng",
) -> Optional[ParsedDocument]:
    """
    Route *meta.abs_path* to the correct parser based on its extension.

    Args:
        meta:               Document metadata (provides abs_path + source_format).
        min_chars_per_page: OCR trigger threshold for PDFs.
        ocr_engine:         "tesseract" | "easyocr".
        ocr_dpi:            Render resolution for OCR.
        ocr_language:       Tesseract language code.

    Returns:
        :class:`ParsedDocument` on success, ``None`` on unrecoverable failure.
    """
    path = meta.abs_path
    fmt = meta.source_format

    log.debug("dispatch_parser", path=str(path), format=fmt)

    try:
        if fmt == "pdf":
            return parse_pdf(
                path, meta, min_chars_per_page, ocr_engine, ocr_dpi, ocr_language
            )
        if fmt == "pptx":
            return parse_pptx(path)
        if fmt == "docx":
            return parse_docx(path)
        if fmt in ("md", "txt"):
            return parse_text(path)

        log.error("dispatch_unknown_format", path=str(path), format=fmt)
        return None

    except Exception as exc:
        log.error(
            "dispatch_parser_exception",
            path=str(path),
            format=fmt,
            exc_type=type(exc).__name__,
            error=str(exc)[:300],
        )
        return None
